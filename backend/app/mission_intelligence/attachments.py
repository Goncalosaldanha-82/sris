from __future__ import annotations

import base64
import hashlib
import io
import json
import mimetypes
import os
import re
import warnings
import zipfile
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path
from uuid import uuid4

from cryptography.hazmat.primitives import hashes
from cryptography.hazmat.primitives.ciphers.aead import AESGCM
from cryptography.hazmat.primitives.kdf.hkdf import HKDF
from sqlalchemy import func
from sqlalchemy.orm import Session

from app.atlas_platform.audit import record_audit
from app.atlas_platform.config import settings

from .models import (
    CanonicalMission,
    MissionArchiveChunk,
    MissionAttachment,
    MissionDialogueSession,
    MissionDialogueTurn,
)
from .mission_archive import index_attachment_text


MAX_ATTACHMENT_BYTES = 20 * 1024 * 1024
MAX_ARCHIVE_ENTRIES = 5_000
MAX_ARCHIVE_UNCOMPRESSED_BYTES = 150 * 1024 * 1024
ATTACHMENT_QUERY_BATCH_SIZE = 400

ALLOWED_EXTENSIONS = {
    ".pdf",
    ".html",
    ".htm",
    ".md",
    ".markdown",
    ".txt",
    ".csv",
    ".tsv",
    ".xlsx",
    ".xls",
    ".docx",
    ".pptx",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".gif",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif"}
OFFICE_ZIP_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
PLAIN_TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".csv", ".tsv"}


class AttachmentError(ValueError):
    def __init__(self, code: str, message: str):
        super().__init__(message)
        self.code = code
        self.message = message


@dataclass(frozen=True)
class PreparedAttachment:
    id: str
    filename: str
    media_type: str
    extension: str
    byte_size: int
    sha256: str
    question_id: str | None
    extracted_text: str
    content: bytes

    @property
    def is_image(self) -> bool:
        return self.extension in IMAGE_EXTENSIONS

    @property
    def is_pdf(self) -> bool:
        return self.extension == ".pdf"


class _VisibleHTML(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self.hidden_depth = 0

    def handle_starttag(self, tag: str, _attrs) -> None:
        if tag.casefold() in {"script", "style", "noscript", "template"}:
            self.hidden_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag.casefold() in {"script", "style", "noscript", "template"}:
            self.hidden_depth = max(0, self.hidden_depth - 1)

    def handle_data(self, data: str) -> None:
        if not self.hidden_depth and data.strip():
            self.parts.append(data.strip())


def _attachment_key(organization_id: str) -> bytes:
    return HKDF(
        algorithm=hashes.SHA256(),
        length=32,
        salt=organization_id.encode("utf-8"),
        info=b"sris-mi-attachment-v1",
    ).derive(settings.jwt_secret.encode("utf-8"))


def _encrypt(
    content: bytes,
    *,
    organization_id: str,
    attachment_id: str,
    purpose: str,
) -> bytes:
    nonce = os.urandom(12)
    aad = f"{organization_id}:{attachment_id}:{purpose}".encode("utf-8")
    return nonce + AESGCM(_attachment_key(organization_id)).encrypt(nonce, content, aad)


def _decrypt_blob(
    encrypted: bytes,
    *,
    organization_id: str,
    attachment_id: str,
    purpose: str,
) -> bytes:
    nonce, ciphertext = encrypted[:12], encrypted[12:]
    aad = f"{organization_id}:{attachment_id}:{purpose}".encode("utf-8")
    try:
        return AESGCM(_attachment_key(organization_id)).decrypt(nonce, ciphertext, aad)
    except Exception as exc:
        raise AttachmentError(
            "attachment_integrity_failed",
            "O anexo falhou a verificação de integridade.",
        ) from exc


def _decrypt(row: MissionAttachment) -> bytes:
    return _decrypt_blob(
        row.encrypted_content,
        organization_id=row.organization_id,
        attachment_id=row.id,
        purpose="content",
    )


def _decrypt_extracted_text(row: MissionAttachment) -> str:
    if not row.extracted_text:
        return ""
    try:
        encrypted = base64.b64decode(row.extracted_text.encode("ascii"), validate=True)
        return _decrypt_blob(
            encrypted,
            organization_id=row.organization_id,
            attachment_id=row.id,
            purpose="extracted-text",
        ).decode("utf-8")
    except AttachmentError:
        raise
    except Exception as exc:
        raise AttachmentError(
            "attachment_integrity_failed",
            "O texto extraído do anexo falhou a verificação de integridade.",
        ) from exc


def _decode_text(content: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-16", "cp1252"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


def _normalise_text(value: str) -> str:
    value = value.replace("\x00", " ")
    value = re.sub(r"[\t\r ]+", " ", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def _extract_zip_xml(content: bytes, extension: str) -> str:
    prefixes = {
        ".docx": ("word/document.xml", "word/header", "word/footer", "word/footnotes.xml"),
        ".pptx": ("ppt/slides/slide", "ppt/notesSlides/notesSlide"),
        ".xlsx": ("xl/sharedStrings.xml", "xl/worksheets/sheet"),
    }[extension]
    parts: list[str] = []
    with zipfile.ZipFile(io.BytesIO(content)) as archive:
        for name in sorted(archive.namelist()):
            if not name.endswith(".xml") or not name.startswith(prefixes):
                continue
            raw = archive.read(name).decode("utf-8", errors="ignore")
            text = re.sub(r"<[^>]+>", " ", raw)
            text = re.sub(r"\s+", " ", text).strip()
            if text:
                parts.append(text)
    return "\n".join(parts)


def _extract_pdf(content: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        return ""
    reader = PdfReader(io.BytesIO(content))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def _extract_docx(content: bytes) -> str:
    try:
        from docx import Document
    except ImportError:
        return _extract_zip_xml(content, ".docx")
    document = Document(io.BytesIO(content))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _extract_xlsx(content: bytes) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return _extract_zip_xml(content, ".xlsx")
    workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"[Folha: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) for value in row if value is not None]
            if values:
                parts.append(" | ".join(values))
    return "\n".join(parts)


def _extract_xls(content: bytes) -> str:
    try:
        import xlrd
    except ImportError:
        return ""
    workbook = xlrd.open_workbook(file_contents=content, on_demand=True)
    parts: list[str] = []
    for sheet in workbook.sheets():
        parts.append(f"[Folha: {sheet.name}]")
        for row_index in range(sheet.nrows):
            values = [str(value) for value in sheet.row_values(row_index) if value not in (None, "")]
            if values:
                parts.append(" | ".join(values))
    return "\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return _extract_zip_xml(content, ".pptx")
    presentation = Presentation(io.BytesIO(content))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"[Diapositivo {index}]")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text.strip())
    return "\n".join(parts)


def _extract_text(content: bytes, extension: str) -> tuple[str, str, str]:
    try:
        if extension in PLAIN_TEXT_EXTENSIONS:
            text = _decode_text(content)
        elif extension in {".html", ".htm"}:
            parser = _VisibleHTML()
            parser.feed(_decode_text(content))
            text = "\n".join(parser.parts)
        elif extension == ".pdf":
            text = _extract_pdf(content)
        elif extension == ".docx":
            text = _extract_docx(content)
        elif extension == ".xlsx":
            text = _extract_xlsx(content)
        elif extension == ".xls":
            text = _extract_xls(content)
        elif extension == ".pptx":
            text = _extract_pptx(content)
        else:
            return "", "visual_ready", ""
        text = _normalise_text(text)
        if text:
            return text, "ready", ""
        if extension == ".pdf":
            return "", "provider_ready", "Texto local indisponível; o PDF será lido visualmente pelo modelo."
        if extension == ".xls":
            return "", "provider_ready", "A folha XLS será enviada ao modelo como ficheiro."
        return "", "partial", "O ficheiro não contém texto extraível."
    except Exception as exc:
        if extension in {".pdf", ".xls"}:
            return "", "provider_ready", f"Extração local indisponível ({type(exc).__name__}); leitura pelo modelo ativada."
        raise AttachmentError(
            "attachment_extraction_failed",
            "Não foi possível ler o conteúdo deste ficheiro.",
        ) from exc


def _safe_filename(filename: str) -> tuple[str, str]:
    clean = (
        Path((filename or "anexo").replace("\\", "/"))
        .name.strip()
        .replace("\x00", "")
    )
    if not clean or len(clean) > 500:
        raise AttachmentError("invalid_filename", "O nome do ficheiro não é válido.")
    extension = Path(clean).suffix.casefold()
    if extension not in ALLOWED_EXTENSIONS:
        allowed = ", ".join(sorted(ALLOWED_EXTENSIONS))
        raise AttachmentError(
            "unsupported_file_type",
            f"Formato não suportado. Formatos aceites: {allowed}.",
        )
    return clean, extension


def _validate_signature(content: bytes, extension: str) -> None:
    if extension == ".pdf" and not content.startswith(b"%PDF-"):
        raise AttachmentError("invalid_file_signature", "O ficheiro não é um PDF válido.")
    if extension in OFFICE_ZIP_EXTENSIONS:
        try:
            with zipfile.ZipFile(io.BytesIO(content)) as archive:
                names = archive.namelist()
                required_member = {
                    ".docx": "word/document.xml",
                    ".xlsx": "xl/workbook.xml",
                    ".pptx": "ppt/presentation.xml",
                }[extension]
                if "[Content_Types].xml" not in names or required_member not in names:
                    raise AttachmentError("invalid_file_signature", "O ficheiro Office não é válido.")
                if len(names) > MAX_ARCHIVE_ENTRIES or sum(
                    item.file_size for item in archive.infolist()
                ) > MAX_ARCHIVE_UNCOMPRESSED_BYTES:
                    raise AttachmentError(
                        "unsafe_archive",
                        "O ficheiro Office excede os limites seguros de descompressão.",
                    )
        except zipfile.BadZipFile as exc:
            raise AttachmentError("invalid_file_signature", "O ficheiro Office está danificado.") from exc
    if extension == ".png" and not content.startswith(b"\x89PNG\r\n\x1a\n"):
        raise AttachmentError("invalid_file_signature", "A imagem PNG não é válida.")
    if extension in {".jpg", ".jpeg"} and not content.startswith(b"\xff\xd8\xff"):
        raise AttachmentError("invalid_file_signature", "A imagem JPEG não é válida.")
    if extension == ".gif" and not content.startswith((b"GIF87a", b"GIF89a")):
        raise AttachmentError("invalid_file_signature", "A imagem GIF não é válida.")
    if extension == ".webp" and not (
        content.startswith(b"RIFF") and content[8:12] == b"WEBP"
    ):
        raise AttachmentError("invalid_file_signature", "A imagem WebP não é válida.")
    if extension in IMAGE_EXTENSIONS:
        try:
            from PIL import Image
        except ImportError:
            return
        expected_format = {
            ".png": "PNG",
            ".jpg": "JPEG",
            ".jpeg": "JPEG",
            ".webp": "WEBP",
            ".gif": "GIF",
        }[extension]
        try:
            with warnings.catch_warnings():
                warnings.simplefilter("error", Image.DecompressionBombWarning)
                with Image.open(io.BytesIO(content)) as image:
                    if image.format != expected_format:
                        raise AttachmentError(
                            "invalid_file_signature",
                            "O formato real da imagem não corresponde ao nome do ficheiro.",
                        )
                    image.verify()
        except AttachmentError:
            raise
        except Exception as exc:
            raise AttachmentError(
                "invalid_file_signature",
                "A imagem está danificada ou excede os limites seguros.",
            ) from exc


def _mission_or_error(db: Session, *, organization_id: str, mission_code: str) -> CanonicalMission:
    row = (
        db.query(CanonicalMission)
        .filter(
            CanonicalMission.organization_id == organization_id,
            CanonicalMission.code == mission_code,
        )
        .one_or_none()
    )
    if row is None:
        raise AttachmentError(
            "mission_not_found",
            "Inicie primeiro a Mission Intelligence para criar o registo canónico desta missão.",
        )
    return row


def attachment_chunk_counts(
    db: Session,
    rows: list[MissionAttachment],
) -> dict[str, int]:
    """Return exact derived-index counts without an N+1 relationship load."""

    attachment_ids = [row.id for row in rows]
    counts: dict[str, int] = {}
    for offset in range(0, len(attachment_ids), ATTACHMENT_QUERY_BATCH_SIZE):
        batch = attachment_ids[offset : offset + ATTACHMENT_QUERY_BATCH_SIZE]
        counts.update(
            {
                str(attachment_id): int(chunk_count)
                for attachment_id, chunk_count in (
                    db.query(
                        MissionArchiveChunk.attachment_id,
                        func.count(MissionArchiveChunk.id),
                    )
                    .filter(MissionArchiveChunk.attachment_id.in_(batch))
                    .group_by(MissionArchiveChunk.attachment_id)
                    .all()
                )
                if attachment_id is not None
            }
        )
    return counts


def attachment_view(
    row: MissionAttachment,
    *,
    archive_chunk_count: int = 0,
) -> dict:
    archive_chunk_count = max(0, int(archive_chunk_count))
    return {
        "id": row.id,
        "evidence_id": f"ATT-{row.id[:8].upper()}",
        "mission_code": row.mission_code,
        "dialogue_session_id": row.dialogue_session_id,
        "question_id": row.question_id,
        "filename": row.original_filename,
        "media_type": row.media_type,
        "extension": row.extension,
        "byte_size": row.byte_size,
        "sha256": row.sha256,
        "extraction_status": row.extraction_status,
        "extraction_error": row.extraction_error,
        "archive_indexed": archive_chunk_count > 0,
        "archive_chunk_count": archive_chunk_count,
        "created_at": row.created_at,
        "epistemic_status": "user_supplied_source",
        "verification_status": "in_review",
    }


def attachment_views(
    db: Session,
    rows: list[MissionAttachment],
) -> list[dict]:
    counts = attachment_chunk_counts(db, rows)
    return [
        attachment_view(
            row,
            archive_chunk_count=counts.get(row.id, 0),
        )
        for row in rows
    ]


def create_attachment(
    db: Session,
    *,
    organization_id: str,
    mission_code: str,
    user_id: str,
    filename: str,
    declared_media_type: str | None,
    content: bytes,
    dialogue_session_id: str | None = None,
    question_id: str | None = None,
) -> MissionAttachment:
    mission = _mission_or_error(db, organization_id=organization_id, mission_code=mission_code)
    clean_name, extension = _safe_filename(filename)
    if not content:
        raise AttachmentError("empty_file", "O ficheiro está vazio.")
    if len(content) > MAX_ATTACHMENT_BYTES:
        raise AttachmentError("file_too_large", "Cada ficheiro pode ter no máximo 20 MB.")
    _validate_signature(content, extension)

    if dialogue_session_id:
        valid_session = (
            db.query(MissionDialogueSession.id)
            .filter(
                MissionDialogueSession.id == dialogue_session_id,
                MissionDialogueSession.organization_id == organization_id,
                MissionDialogueSession.mission_id == mission.id,
            )
            .scalar()
        )
        if valid_session is None:
            raise AttachmentError("session_not_found", "A sessão indicada não pertence a esta missão.")

    digest = hashlib.sha256(content).hexdigest()
    existing = (
        db.query(MissionAttachment)
        .filter(
            MissionAttachment.organization_id == organization_id,
            MissionAttachment.mission_id == mission.id,
            MissionAttachment.sha256 == digest,
        )
        .one_or_none()
    )
    if existing is not None:
        return existing

    extracted_text, extraction_status, extraction_error = _extract_text(content, extension)
    attachment_id = str(uuid4())
    clean_question_id = (question_id or "").strip()
    if len(clean_question_id) > 120:
        raise AttachmentError(
            "invalid_question_id",
            "O identificador da pergunta não é válido.",
        )
    declared_type = (
        (declared_media_type or "").split(";", 1)[0].strip().casefold()[:160]
    )
    media_type = (
        mimetypes.guess_type(clean_name)[0]
        or declared_type
        or "application/octet-stream"
    )
    row = MissionAttachment(
        id=attachment_id,
        organization_id=organization_id,
        mission_id=mission.id,
        mission_code=mission.code,
        dialogue_session_id=dialogue_session_id,
        question_id=clean_question_id or None,
        original_filename=clean_name,
        media_type=media_type,
        extension=extension,
        byte_size=len(content),
        sha256=digest,
        encrypted_content=_encrypt(
            content,
            organization_id=organization_id,
            attachment_id=attachment_id,
            purpose="content",
        ),
        extracted_text=(
            base64.b64encode(
                _encrypt(
                    extracted_text.encode("utf-8"),
                    organization_id=organization_id,
                    attachment_id=attachment_id,
                    purpose="extracted-text",
                )
            ).decode("ascii")
            if extracted_text
            else ""
        ),
        extraction_status=extraction_status,
        extraction_error=extraction_error,
        created_by_user_id=user_id,
    )
    db.add(row)
    if extracted_text:
        index_attachment_text(
            db,
            attachment=row,
            extracted_text=extracted_text,
        )
    record_audit(
        db,
        action="mission_intelligence.attachment_uploaded",
        resource_type="mission_attachment",
        resource_id=row.id,
        organization_id=organization_id,
        user_id=user_id,
        payload={
            "mission_code": mission.code,
            "filename": clean_name,
            "byte_size": len(content),
            "sha256": digest,
            "question_id": row.question_id,
            "archive_indexed": bool(extracted_text),
        },
    )
    db.commit()
    db.refresh(row)
    return row


def list_attachments(db: Session, *, organization_id: str, mission_code: str) -> list[MissionAttachment]:
    return (
        db.query(MissionAttachment)
        .filter(
            MissionAttachment.organization_id == organization_id,
            MissionAttachment.mission_code == mission_code,
        )
        .order_by(MissionAttachment.created_at.asc())
        .all()
    )


def get_attachment(
    db: Session,
    *,
    organization_id: str,
    mission_code: str,
    attachment_id: str,
) -> MissionAttachment | None:
    return (
        db.query(MissionAttachment)
        .filter(
            MissionAttachment.id == attachment_id,
            MissionAttachment.organization_id == organization_id,
            MissionAttachment.mission_code == mission_code,
        )
        .one_or_none()
    )


def attachment_content(row: MissionAttachment) -> bytes:
    return _decrypt(row)


def prepare_turn_attachments(
    db: Session,
    *,
    organization_id: str,
    mission_id: str,
    attachment_ids: list[str],
) -> list[PreparedAttachment]:
    rows = prepare_turn_attachment_rows(
        db,
        organization_id=organization_id,
        mission_id=mission_id,
        attachment_ids=attachment_ids,
    )
    return [
        PreparedAttachment(
            id=row.id,
            filename=row.original_filename,
            media_type=row.media_type,
            extension=row.extension,
            byte_size=row.byte_size,
            sha256=row.sha256,
            question_id=row.question_id,
            extracted_text=_decrypt_extracted_text(row),
            content=_decrypt(row),
        )
        for row in rows
    ]


def prepare_turn_attachment_rows(
    db: Session,
    *,
    organization_id: str,
    mission_id: str,
    attachment_ids: list[str],
) -> list[MissionAttachment]:
    """Resolve an arbitrarily large ID set in database-safe batches."""

    if not attachment_ids:
        return []
    by_id: dict[str, MissionAttachment] = {}
    for offset in range(0, len(attachment_ids), ATTACHMENT_QUERY_BATCH_SIZE):
        batch = attachment_ids[offset : offset + ATTACHMENT_QUERY_BATCH_SIZE]
        rows = (
            db.query(MissionAttachment)
            .filter(
                MissionAttachment.organization_id == organization_id,
                MissionAttachment.mission_id == mission_id,
                MissionAttachment.id.in_(batch),
            )
            .all()
        )
        by_id.update((row.id, row) for row in rows)
    missing = [item for item in attachment_ids if item not in by_id]
    if missing:
        raise AttachmentError(
            "attachment_not_found",
            "Um ou mais anexos não pertencem a esta missão.",
        )
    return [by_id[item] for item in attachment_ids]


def backfill_mission_archive_index(
    db: Session,
    *,
    organization_id: str,
    mission_id: str,
    priority_attachment_ids: list[str] | None = None,
    batch_size: int = 64,
) -> int:
    """Lazily index legacy encrypted attachments without blocking the archive.

    New uploads are indexed synchronously. This bounded backfill only exists
    for sources created before the scalable archive migration.
    """

    priority = list(dict.fromkeys(priority_attachment_ids or []))
    rows: list[MissionAttachment] = []
    if priority:
        for offset in range(0, len(priority), ATTACHMENT_QUERY_BATCH_SIZE):
            batch = priority[offset : offset + ATTACHMENT_QUERY_BATCH_SIZE]
            rows.extend(
                db.query(MissionAttachment)
                .filter(
                    MissionAttachment.organization_id == organization_id,
                    MissionAttachment.mission_id == mission_id,
                    MissionAttachment.id.in_(batch),
                    MissionAttachment.extracted_text != "",
                    ~MissionAttachment.archive_chunks.any(),
                )
                .limit(max(0, batch_size - len(rows)))
                .all()
            )
            if len(rows) >= batch_size:
                break
    remaining = max(0, batch_size - len(rows))
    if remaining:
        priority_set = {row.id for row in rows}
        query = db.query(MissionAttachment).filter(
            MissionAttachment.organization_id == organization_id,
            MissionAttachment.mission_id == mission_id,
            MissionAttachment.extracted_text != "",
            ~MissionAttachment.archive_chunks.any(),
        )
        if priority_set:
            query = query.filter(~MissionAttachment.id.in_(priority_set))
        rows.extend(query.order_by(MissionAttachment.created_at.asc()).limit(remaining).all())

    indexed = 0
    for row in rows:
        text = _decrypt_extracted_text(row)
        if text:
            index_attachment_text(db, attachment=row, extracted_text=text)
            indexed += 1
    if indexed:
        db.flush()
    return indexed


def delete_attachment(
    db: Session,
    *,
    row: MissionAttachment,
    user_id: str,
) -> None:
    referenced_turns = (
        db.query(MissionDialogueTurn.attachment_ids_json)
        .join(MissionDialogueSession)
        .filter(MissionDialogueSession.mission_id == row.mission_id)
        .all()
    )
    is_referenced = False
    for (value,) in referenced_turns:
        try:
            attachment_ids = json.loads(value or "[]")
        except (TypeError, ValueError):
            attachment_ids = []
        if row.id in attachment_ids:
            is_referenced = True
            break
    if is_referenced:
        raise AttachmentError(
            "attachment_in_use",
            "Este anexo já faz parte do histórico auditável e não pode ser eliminado.",
        )
    record_audit(
        db,
        action="mission_intelligence.attachment_deleted",
        resource_type="mission_attachment",
        resource_id=row.id,
        organization_id=row.organization_id,
        user_id=user_id,
        payload={"mission_code": row.mission_code, "filename": row.original_filename, "sha256": row.sha256},
    )
    db.delete(row)
    db.commit()
